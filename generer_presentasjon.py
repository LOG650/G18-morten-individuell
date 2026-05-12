"""
generer_presentasjon.py
Genererer PowerPoint-presentasjon for muntlig eksamen LOG650.
Kjør: python generer_presentasjon.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUTPUT = "Presentasjon_LOG650_ME_Morten_Eidsvag.pptx"

# --- Fargepalett ---
NAVY   = RGBColor(0x1F, 0x38, 0x64)
BLUE   = RGBColor(0x2E, 0x75, 0xB6)
LBLUE  = RGBColor(0xBD, 0xD7, 0xEE)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x59, 0x59, 0x59)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE = RGBColor(0xC5, 0x5A, 0x11)
GREEN  = RGBColor(0x37, 0x5F, 0x1B)
RED    = RGBColor(0xC0, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

W = prs.slide_width
H = prs.slide_height

# ── Hjelpefunksjoner ──────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=18, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def title_bar(slide, title):
    rect(slide, 0, 0, W, Inches(1.15), fill=NAVY)
    rect(slide, 0, Inches(1.15), W, Inches(0.04), fill=BLUE)
    txt(slide, title, Inches(0.45), Inches(0.15), Inches(12.4), Inches(0.9),
        size=26, bold=True, color=WHITE)

def footer(slide, page_num):
    rect(slide, 0, Inches(7.2), W, Inches(0.3), fill=NAVY)
    txt(slide, "Morten Eidsvåg  |  LOG650  |  Høgskolen i Molde  |  Mai 2026",
        Inches(0.3), Inches(7.2), Inches(11), Inches(0.3),
        size=9, color=WHITE)
    txt(slide, str(page_num), Inches(12.8), Inches(7.2), Inches(0.5), Inches(0.3),
        size=9, color=WHITE, align=PP_ALIGN.RIGHT)

def bullets(slide, items, l, t, w, h, size=17, color=BLACK, indent=0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.level = item.get("level", indent)
        p.space_before = Pt(item.get("space", 6))
        r = p.add_run()
        r.text = item["text"]
        r.font.size = Pt(item.get("size", size))
        r.font.bold = item.get("bold", False)
        r.font.italic = item.get("italic", False)
        r.font.color.rgb = item.get("color", color)
    return tb

def kpi_box(slide, value, label, l, t, w=Inches(2.6), h=Inches(1.3),
            bg=BLUE, val_color=WHITE, lbl_color=WHITE):
    rect(slide, l, t, w, h, fill=bg)
    txt(slide, value, l, t + Inches(0.1), w, Inches(0.75),
        size=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, l, t + Inches(0.8), w, Inches(0.45),
        size=12, color=lbl_color, align=PP_ALIGN.CENTER)

def add_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text

def tbl(slide, data, headers, l, t, w, h,
        col_widths=None, header_bg=NAVY, row_alt=LGRAY,
        font_size=14, header_font_size=13):
    rows = len(data) + 1
    cols = len(headers)
    table = slide.shapes.add_table(rows, cols, l, t, w, h).table
    col_w = w // cols
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw

    # Header
    for j, h_txt in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h_txt
        r.font.size = Pt(header_font_size)
        r.font.bold = True; r.font.color.rgb = WHITE

    # Data rows
    for i, row in enumerate(data):
        bg = row_alt if i % 2 == 1 else WHITE
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font_size)
            if isinstance(val, str) and ("XGBoost" in val or "52,19" in val or "811" in val or "702" in val):
                r.font.bold = True; r.font.color.rgb = GREEN
    return table

# ── SLIDE 1: Tittel ───────────────────────────────────────────────────────────
s = new_slide()
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, Inches(4.8), W, Inches(0.06), fill=BLUE)

txt(s, "Etterspørselsprognoser for\nfarmasøytiske varelinjer",
    Inches(0.8), Inches(0.8), Inches(11.7), Inches(2.4),
    size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Sammenligning av naiv sesongprognose, Holt-Winters,\nARIMA/SARIMA og XGBoost",
    Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.2),
    size=20, color=LBLUE, align=PP_ALIGN.CENTER, italic=True)

txt(s, "Morten Eidsvåg",
    Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.5),
    size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "LOG650 – Forskningsprosjekt: Logistikk og kunstig intelligens\nHøgskolen i Molde  |  Mai 2026",
    Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.8),
    size=14, color=LBLUE, align=PP_ALIGN.CENTER)

add_notes(s, "Introduser deg selv og temaet. Presenter problemet kort: økt legemiddelmangel, behov for bedre prognoser, og at du sammenligner fire ulike modeller på reelle data.")

# ── SLIDE 2: Agenda ───────────────────────────────────────────────────────────
s = new_slide()
title_bar(s, "Agenda")
footer(s, 2)

agenda = [
    {"text": "1   Bakgrunn og motivasjon",   "size": 19, "space": 10},
    {"text": "2   Problemstilling og forskningsspørsmål", "size": 19, "space": 10},
    {"text": "3   Case og datasett",          "size": 19, "space": 10},
    {"text": "4   Metode og modeller",        "size": 19, "space": 10},
    {"text": "5   Resultater",                "size": 19, "space": 10},
    {"text": "6   Diskusjon",                 "size": 19, "space": 10},
    {"text": "7   Konklusjon og anbefalinger","size": 19, "space": 10},
]
bullets(s, agenda, Inches(2.5), Inches(1.5), Inches(8.5), Inches(5.5))

add_notes(s, "Gi en rask oversikt over presentasjonsstrukturen. Ca. 30 minutter totalt.")

# ── SLIDE 3: Bakgrunn ─────────────────────────────────────────────────────────
s = new_slide()
title_bar(s, "Bakgrunn og motivasjon")
footer(s, 3)

kpi_box(s, "1 403", "mangelrapporter\n(Statens legemiddelverk, 2023)",
        Inches(0.4), Inches(1.4), w=Inches(2.8), h=Inches(1.5))
kpi_box(s, "76 %", "økning\nfra 2022 til 2023",
        Inches(0.4), Inches(3.1), w=Inches(2.8), h=Inches(1.5), bg=ORANGE)

items = [
    {"text": "Feilprognoser har to konsekvenser:", "bold": True, "space": 8, "size": 17},
    {"text": "  →  For lav: legemiddelmangel — tomme apotekhyller", "size": 16, "space": 4},
    {"text": "  →  For høy: overlagring, kapitalbinding, ukurans", "size": 16, "space": 4},
    {"text": "", "size": 8, "space": 2},
    {"text": "Tradisjonelle metoder: Holt-Winters og ARIMA", "bold": True, "space": 8, "size": 17},
    {"text": "  →  Statistiske tidsseriemodeller, veletablerte", "size": 16, "space": 4},
    {"text": "", "size": 8, "space": 2},
    {"text": "Nyere metoder: maskinlæring (XGBoost)", "bold": True, "space": 8, "size": 17},
    {"text": "  →  Kan håndtere komplekse, ikke-lineære mønstre", "size": 16, "space": 4},
    {"text": "", "size": 8, "space": 2},
    {"text": "Gap i litteraturen:", "bold": True, "space": 8, "size": 17, "color": BLUE},
    {"text": "  →  Få studier på paneldata med mange SKU-er,", "size": 16, "space": 4},
    {"text": "       kort historikk og i et distribusjonsledd", "size": 16, "space": 4},
]
bullets(s, items, Inches(3.7), Inches(1.35), Inches(9.2), Inches(5.8))

add_notes(s, "Presenter problemet: Norge har hatt økt legemiddelmangel. Feilprognoser er en bidragsyter. Tradisjonelle og nyere metoder finnes, men ingen har systematisk sammenlignet dem på nettopp denne typen datasett.")

# ── SLIDE 4: Problemstilling ──────────────────────────────────────────────────
s = new_slide()
title_bar(s, "Problemstilling og forskningsspørsmål")
footer(s, 4)

rect(s, Inches(0.4), Inches(1.35), Inches(12.5), Inches(1.5), fill=LBLUE)
txt(s,
    "«Hvilke prognosemodeller egner seg best for månedlige leveranser av farmasøytiske "
    "varelinjer i et sentrallager–grossist-system, og i hvilken grad varierer "
    "modellprestasjonene på tvers av produktsegmenter?»",
    Inches(0.6), Inches(1.42), Inches(12.1), Inches(1.4),
    size=16, bold=True, color=NAVY, italic=True)

items = [
    {"text": "RQ1 — Hvilken modell gir lavest prognosefeil (RMSE, MAE, MAPE)\n"
             "          målt på tvers av alle 105 SKU-er?",
     "bold": True, "size": 16, "space": 14, "color": NAVY},
    {"text": "RQ2 — Varierer den relative modellprestasjonen mellom\n"
             "          produktsegmenter (volum og sesongmønster)?",
     "bold": True, "size": 16, "space": 14, "color": NAVY},
    {"text": "RQ3 — I hvilken grad gir økt modellkompleksitet reell merverdi\n"
             "          i en operasjonell planleggingskontekst?",
     "bold": True, "size": 16, "space": 14, "color": NAVY},
]
bullets(s, items, Inches(0.5), Inches(3.1), Inches(12.4), Inches(3.8))

add_notes(s, "Les problemstillingen høyt. Forklar at den er todelt: (1) hvilken modell er best samlet, og (2) om svaret er likt for alle produkttyper. De tre RQ-ene operasjonaliserer problemstillingen.")

# ── SLIDE 5: Case og datasett ─────────────────────────────────────────────────
s = new_slide()
title_bar(s, "Case og datasett")
footer(s, 5)

# Supply chain boxes
for i, (lbl, x) in enumerate([
    ("Europeisk\nsentrallager", 0.4),
    ("Norske\ngrossister",      3.2),
    ("Apotek",                  6.0),
    ("Pasient",                 8.8),
]):
    rect(s, Inches(x), Inches(1.4), Inches(2.4), Inches(0.9),
         fill=NAVY if i == 0 else (BLUE if i == 1 else LGRAY),
         line=NAVY)
    txt(s, lbl,
        Inches(x), Inches(1.4), Inches(2.4), Inches(0.9),
        size=13, bold=(i<2), color=WHITE if i<2 else BLACK,
        align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s, "→", Inches(x+2.4), Inches(1.5), Inches(0.4), Inches(0.7),
            size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

rect(s, Inches(0.4), Inches(2.45), Inches(5.2), Inches(0.35), fill=ORANGE)
txt(s, "← Denne studien analyserer dette leddet",
    Inches(0.5), Inches(2.47), Inches(5.0), Inches(0.3),
    size=12, color=WHITE, bold=True)

# KPI boxes
for val, lbl, x in [
    ("105", "SKU-er", 0.4),
    ("48", "måneder", 3.2),
    ("76 %", "har nullverdier", 6.0),
    ("5 040", "observasjoner", 9.4),
]:
    kpi_box(s, val, lbl, Inches(x), Inches(3.05), w=Inches(2.5), h=Inches(1.2))

txt(s, "Treningssett: Jan 2022 – Des 2024  (36 måneder)",
    Inches(0.4), Inches(4.45), Inches(5.8), Inches(0.45),
    size=15, bold=True, color=NAVY)
txt(s, "Testsett: Jan – Des 2025  (12 måneder)",
    Inches(0.4), Inches(4.9), Inches(5.8), Inches(0.45),
    size=15, bold=True, color=ORANGE)
txt(s, "Anonymisert datasett fra virksomhetens ERP-system",
    Inches(6.5), Inches(4.45), Inches(6.4), Inches(0.45),
    size=14, italic=True, color=GRAY)

add_notes(s, "Beskriv forsyningskjedestrukturen. Understreke at prognosene gjelder første ledd — fra sentrallager til grossist. Nevn at 76% av SKU-ene har nullverdier, noe som påvirker MAPE-beregningen.")

# ── SLIDE 6: Metode ───────────────────────────────────────────────────────────
s = new_slide()
title_bar(s, "Metode og forskningsdesign")
footer(s, 6)

left_items = [
    {"text": "Kvantitativ casestudie", "bold": True, "size": 17, "space": 10},
    {"text": "  Positivistisk perspektiv — virkeligheten kan måles", "size": 15, "space": 3},
    {"text": "  Ett datasett, én virksomhet, full metodisk kontroll", "size": 15, "space": 3},
    {"text": "", "size": 8},
    {"text": "Hold-out-validering", "bold": True, "size": 17, "space": 10},
    {"text": "  36 mnd trening  →  12 mnd test (out-of-sample)", "size": 15, "space": 3},
    {"text": "  Testperioden dekker én full sesongssyklus", "size": 15, "space": 3},
    {"text": "", "size": 8},
    {"text": "Felles evalueringsrammeverk", "bold": True, "size": 17, "space": 10},
    {"text": "  Alle modeller: identisk data, split og feilmål", "size": 15, "space": 3},
]
bullets(s, left_items, Inches(0.4), Inches(1.35), Inches(6.0), Inches(5.8))

right_items = [
    {"text": "Segmentanalyse (RQ2)", "bold": True, "size": 17, "space": 10},
    {"text": "  Sesongstyrke: Fs ≥ 0,64 (sterk / svak)", "size": 15, "space": 3},
    {"text": "  Volumnivå: tertiler (høy / middels / lav)", "size": 15, "space": 3},
    {"text": "", "size": 8},
    {"text": "Statistisk testing (RQ3)", "bold": True, "size": 17, "space": 10},
    {"text": "  Diebold-Mariano-test (DM)", "size": 15, "space": 3},
    {"text": "  Signifikansnivå α = 0,05", "size": 15, "space": 3},
    {"text": "  To par: HW vs. ARIMA, XGBoost vs. ARIMA", "size": 15, "space": 3},
    {"text": "", "size": 8},
    {"text": "Kontrolleksperiment", "bold": True, "size": 17, "space": 10},
    {"text": "  XGBoost trent individuelt per SKU", "size": 15, "space": 3},
    {"text": "  Isolerer pooling-effekten", "size": 15, "space": 3},
]
bullets(s, right_items, Inches(6.8), Inches(1.35), Inches(6.1), Inches(5.8))

rect(s, Inches(6.6), Inches(1.35), Inches(0.04), Inches(5.7), fill=LBLUE)

add_notes(s, "Forklar designvalget: kvantitativ casestudie for metodisk kontroll. Understreke at hold-out er standard for tidsserier (ikke krysskalibrering — ville gi datalekkasje). DM-testen er viktig for å vurdere om forskjellene er statistisk signifikante.")

# ── SLIDE 7: De fire modellene ────────────────────────────────────────────────
s = new_slide()
title_bar(s, "De fire prognosemodellene")
footer(s, 7)

models = [
    ("Naiv sesongprognose", "Referansemodell (benchmark)",
     "ŷ_{t+h} = y_{t+h−12}",
     "Ingen parametere. Predikerer samme verdi som tilsvarende\nmåned i fjor. Brukes for å vurdere om andre modeller\ngir reell forbedring.", LGRAY, BLACK),
    ("Holt-Winters", "Eksponentiell glatting",
     "Nivå (α)  +  Trend (β)  +  Sesong (γ)",
     "Dekomponerer tidsserien. Additiv eller multiplikativ\nsesong (valg basert på CV > 0,5). Parametere estimeres\nved minimering av SSE.", LBLUE, NAVY),
    ("ARIMA / SARIMA", "Statistisk tidsseriemodell",
     "ARIMA(p,d,q)(P,D,Q)₁₂",
     "Autoregresjon + differensiering + glidende gjennomsnitt.\nParametere valgt automatisk via AIC-minimering\n(auto_arima).", LBLUE, NAVY),
    ("XGBoost", "Global maskinlæringsmodell",
     "Features: lag 1–12, måned, tidsindeks",
     "Ensemble av beslutningstrær. Trent på alle 105 SKU-er\nsamlet (3 780 observasjoner). Hyperparametre via\ntidsserie-krysskalibrering.", NAVY, WHITE),
]

for i, (name, sub, formula, desc, bg, fg) in enumerate(models):
    col = i % 2
    row = i // 2
    x = Inches(0.3 + col * 6.5)
    y = Inches(1.35 + row * 2.9)
    w = Inches(6.3)
    h = Inches(2.7)
    rect(s, x, y, w, h, fill=bg)
    txt(s, name, x + Inches(0.15), y + Inches(0.12),
        w - Inches(0.3), Inches(0.45), size=17, bold=True, color=fg)
    txt(s, sub, x + Inches(0.15), y + Inches(0.5),
        w - Inches(0.3), Inches(0.35), size=12, italic=True,
        color=GRAY if bg != NAVY else LBLUE)
    txt(s, formula, x + Inches(0.15), y + Inches(0.82),
        w - Inches(0.3), Inches(0.4), size=13, bold=True,
        color=BLUE if bg != NAVY else WHITE)
    txt(s, desc, x + Inches(0.15), y + Inches(1.25),
        w - Inches(0.3), Inches(1.3), size=12,
        color=BLACK if bg != NAVY else LBLUE)

add_notes(s, "Gå gjennom hver modell konseptuelt. Understreke at naiv er referansemodell — alle andre bør slå den. XGBoost er ikke en tidsseriemodell og krever feature engineering. Global = trent på alle SKU-er samlet.")

# ── SLIDE 8: Feilmål og segmentering ─────────────────────────────────────────
s = new_slide()
title_bar(s, "Evalueringsmål og segmentering")
footer(s, 8)

for name, formula, desc, x in [
    ("RMSE", "√(1/n · Σ(yₜ − ŷₜ)²)", "Sensitiv for store\navvik / uteliggere", 0.35),
    ("MAE",  "1/n · Σ|yₜ − ŷₜ|",      "Robust, intuitiv —\ngjennomsnittlig avvik", 4.55),
    ("MAPE", "100/n · Σ|yₜ−ŷₜ|/yₜ",   "Sammenlignbar på tvers\nav volumnivåer (%)", 8.75),
]:
    rect(s, Inches(x), Inches(1.35), Inches(3.8), Inches(2.3), fill=NAVY)
    txt(s, name, Inches(x + 0.1), Inches(1.42), Inches(3.6), Inches(0.55),
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, formula, Inches(x + 0.1), Inches(1.95), Inches(3.6), Inches(0.55),
        size=14, color=LBLUE, align=PP_ALIGN.CENTER, italic=True)
    txt(s, desc, Inches(x + 0.1), Inches(2.5), Inches(3.6), Inches(0.7),
        size=13, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Aggregering: Median over alle 105 SKU-er  (ikke gjennomsnitt — robust mot ekstremverdier i lavvolum-segmentet)",
    Inches(0.35), Inches(3.75), Inches(12.6), Inches(0.5),
    size=14, color=ORANGE, bold=True)

rect(s, Inches(0.35), Inches(4.35), Inches(6.1), Inches(2.75), fill=LGRAY)
txt(s, "Sesongstyrke (Fs)", Inches(0.5), Inches(4.42),
    Inches(5.8), Inches(0.4), size=15, bold=True, color=NAVY)
items_s = [
    {"text": "Beregnes fra STL-dekomponering:", "size": 14},
    {"text": "Fs = max(0, 1 − Var(Rₜ)/Var(Sₜ+Rₜ))", "size": 13, "italic": True, "color": BLUE},
    {"text": "Fs ≥ 0,64 → «Sterk sesong»  (18 / 105 SKU-er)", "size": 14, "bold": True, "space": 6},
    {"text": "Fs < 0,64 → «Svak/ingen sesong»  (87 / 105 SKU-er)", "size": 14, "space": 4},
]
bullets(s, items_s, Inches(0.5), Inches(4.82), Inches(5.8), Inches(2.2))

rect(s, Inches(6.65), Inches(4.35), Inches(6.3), Inches(2.75), fill=LGRAY)
txt(s, "Volumnivå (tertiler)", Inches(6.8), Inches(4.42),
    Inches(6.0), Inches(0.4), size=15, bold=True, color=NAVY)
items_v = [
    {"text": "Basert på gjennomsnittlig månedlig etterspørsel:", "size": 14},
    {"text": "Høyt volum  —  øverste tertil  (35 SKU-er)", "size": 14, "bold": True, "space": 6},
    {"text": "Middels volum  —  midtre tertil  (35 SKU-er)", "size": 14, "space": 4},
    {"text": "Lavt volum  —  nederste tertil  (35 SKU-er)", "size": 14, "space": 4},
]
bullets(s, items_v, Inches(6.8), Inches(4.82), Inches(6.0), Inches(2.2))

add_notes(s, "Forklar hvorfor vi bruker tre feilmål — de måler ulike ting. RMSE er sensitiv for uteliggere, MAE er robust, MAPE er relativ og sammenlignbar. Median brukes fordi fordelingen er høyreskjev med noen ekstreme verdier for lavvolum-SKU-er.")

# ── SLIDE 9: Samlede resultater ───────────────────────────────────────────────
s = new_slide()
title_bar(s, "Samlede resultater — alle 105 SKU-er")
footer(s, 9)

txt(s, "Median RMSE, MAE og MAPE  |  Testperiode: Jan–Des 2025",
    Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.4),
    size=13, italic=True, color=GRAY)

headers = ["Modell", "RMSE", "MAE", "MAPE (%)"]
data = [
    ["Naiv sesongprognose",  "1 131,65", "914,58", "67,50"],
    ["Holt-Winters",         "935,84",   "769,90", "62,49"],
    ["XGBoost (individuell)","949,67",   "845,80", "62,61"],
    ["ARIMA / SARIMA",       "846,20",   "820,01", "57,77"],
    ["XGBoost (global) ★",   "811,24",   "702,95", "52,19"],
]
col_widths = [Inches(4.0), Inches(2.2), Inches(2.2), Inches(2.2)]
tbl(s, data, headers, Inches(1.3), Inches(1.75), Inches(10.6), Inches(3.2),
    col_widths=col_widths, font_size=15, header_font_size=14)

items = [
    {"text": "★  XGBoost global er best på alle tre feilmål", "bold": True, "size": 15, "color": GREEN, "space": 10},
    {"text": "    MAPE-forbedring: 15,3 pp over naiv  |  10,3 pp over Holt-Winters", "size": 14, "color": BLACK, "space": 4},
    {"text": "★  XGBoost individuell ≈ Holt-Winters  →  pooling-effekten er avgjørende", "bold": True, "size": 15, "color": BLUE, "space": 10},
]
bullets(s, items, Inches(0.4), Inches(5.1), Inches(12.5), Inches(1.8))

add_notes(s, "Dette er kjerneresultatet. XGBoost global er best på alle tre mål. Understreke at XGBoost individuell (trent per SKU) er nesten identisk med Holt-Winters — det viser at det er pooling, ikke gradient boosting-algoritmen, som gir fordelen. Absoluttallene er høye — MAPE 52% er ikke god presisjon i operativ forstand, men det er beste *relative* resultat på et krevende datasett.")

# ── SLIDE 10: Segmentanalyse sesong ──────────────────────────────────────────
s = new_slide()
title_bar(s, "Segmentanalyse — sesongstyrke")
footer(s, 10)

txt(s, "Median MAPE fordelt på SKU-er med sterk (Fs ≥ 0,64) og svak sesong",
    Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.4), size=13, italic=True, color=GRAY)

headers2 = ["Modell", "Sterk sesong (18 SKU-er)", "Svak/ingen sesong (87 SKU-er)"]
data2 = [
    ["Naiv",          "71,92 %", "67,49 %"],
    ["Holt-Winters",  "65,82 %", "61,01 %"],
    ["ARIMA",         "69,86 %", "56,52 %"],
    ["XGBoost",       "79,13 %", "51,50 %"],
]
col_widths2 = [Inches(3.5), Inches(3.5), Inches(3.5)]
tbl(s, data2, headers2, Inches(0.9), Inches(1.75), Inches(10.6), Inches(2.7),
    col_widths=col_widths2, font_size=15)

rect(s, Inches(0.4), Inches(4.6), Inches(5.9), Inches(2.5), fill=LGRAY)
txt(s, "Sterk sesong — Holt-Winters er best",
    Inches(0.55), Inches(4.67), Inches(5.6), Inches(0.4),
    size=15, bold=True, color=NAVY)
items_l = [
    {"text": "XGBoost svakest: 79,13 % — høyere enn naiv (71,92 %)", "size": 14, "color": RED, "bold": True},
    {"text": "Årsak: kun 18/105 SKU-er har sterk sesong.", "size": 13, "space": 4},
    {"text": "Global modell optimerer for de 87 andre —", "size": 13, "space": 4},
    {"text": "sesongpregede serier «drukner».", "size": 13, "space": 4},
]
bullets(s, items_l, Inches(0.55), Inches(5.1), Inches(5.7), Inches(1.9))

rect(s, Inches(6.8), Inches(4.6), Inches(5.9), Inches(2.5), fill=LGRAY)
txt(s, "Svak sesong — XGBoost klart best",
    Inches(6.95), Inches(4.67), Inches(5.6), Inches(0.4),
    size=15, bold=True, color=GREEN)
items_r = [
    {"text": "XGBoost: 51,50 % — 9,5 pp bedre enn Holt-Winters", "size": 14, "color": GREEN, "bold": True},
    {"text": "Global pooling fungerer best når seriene", "size": 13, "space": 4},
    {"text": "ligner hverandre (ikke-sesongpreget adferd).", "size": 13, "space": 4},
]
bullets(s, items_r, Inches(6.95), Inches(5.1), Inches(5.6), Inches(1.9))

add_notes(s, "Dette svarer på RQ2. Modellprestasjonene varierer tydelig mellom segmenter. XGBoost er svakest for sesongpregede SKU-er (79% MAPE — verre enn naiv!) fordi den globale modellen optimerer for det dominerende ikke-sesongpregede mønsteret. Holt-Winters og ARIMA er mer robuste fordi de tilpasser sesongkomponenten per serie uavhengig av andre SKU-er.")

# ── SLIDE 11: Segmentanalyse volum ────────────────────────────────────────────
s = new_slide()
title_bar(s, "Segmentanalyse — volumnivå")
footer(s, 11)

txt(s, "Median MAPE fordelt på tertiler av gjennomsnittlig månedlig etterspørsel",
    Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.4), size=13, italic=True, color=GRAY)

headers3 = ["Modell", "Høyt volum", "Middels volum", "Lavt volum"]
data3 = [
    ["Naiv",         "58,85 %", "67,50 %", "86,07 %"],
    ["Holt-Winters", "61,01 %", "63,63 %", "62,49 %"],
    ["ARIMA",        "54,91 %", "57,77 %", "61,48 %"],
    ["XGBoost",      "49,46 %", "53,32 %", "66,61 %"],
]
col_widths3 = [Inches(3.2), Inches(2.4), Inches(2.4), Inches(2.4)]
tbl(s, data3, headers3, Inches(0.9), Inches(1.75), Inches(10.6), Inches(2.7),
    col_widths=col_widths3, font_size=15)

items = [
    {"text": "Høyt volum:  XGBoost best med 49,46 % — regulære serier gir stabilt laggsignal", "size": 15, "color": GREEN, "bold": True, "space": 12},
    {"text": "Lavt volum:  XGBoost svakest med 66,61 % — støy dominerer, pooling gir skjevhet mot porteføljesnittet", "size": 15, "color": RED, "bold": True, "space": 10},
    {"text": "Holt-Winters og ARIMA: relativt stabil prestasjon på tvers av alle volumnivåer", "size": 14, "color": BLUE, "space": 10},
]
bullets(s, items, Inches(0.4), Inches(4.6), Inches(12.5), Inches(2.5))

add_notes(s, "For høyvolum-SKU-er er etterspørselen mer stabil og laggsignalet sterkere — XGBoost utnytter dette godt. For lavvolum-SKU-er dominerer støy og nullobservasjoner. Den globale modellen trekker prediksjoner mot porteføljesnittet — gunstig for høyvolum-serier som ligner snittet, men et systematisk problem for lavvolum-serier som avviker strukturelt.")

# ── SLIDE 12: Statistisk testing ─────────────────────────────────────────────
s = new_slide()
title_bar(s, "Statistisk testing og kontrolleksperiment")
footer(s, 12)

rect(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(3.5), fill=LGRAY)
txt(s, "Diebold-Mariano-test (DM)",
    Inches(0.5), Inches(1.42), Inches(5.7), Inches(0.45),
    size=16, bold=True, color=NAVY)
txt(s, "H₀: Begge modeller har lik forventet prognosepresisjon  |  α = 0,05",
    Inches(0.5), Inches(1.88), Inches(5.7), Inches(0.45),
    size=12, italic=True, color=GRAY)

headers4 = ["Par", "DM", "p-verdi"]
data4 = [
    ["Holt-Winters vs. ARIMA", "1,766", "0,078"],
    ["XGBoost vs. ARIMA",      "−0,779","0,436"],
]
col_w4 = [Inches(2.8), Inches(1.2), Inches(1.4)]
tbl(s, data4, headers4, Inches(0.5), Inches(2.4), Inches(5.5), Inches(1.6),
    col_widths=col_w4, font_size=14)

txt(s, "Ingen av parene er statistisk signifikant\nforskjellige  (begge p > 0,05)",
    Inches(0.5), Inches(4.1), Inches(5.7), Inches(0.7),
    size=14, bold=True, color=ORANGE)

rect(s, Inches(6.6), Inches(1.35), Inches(6.35), Inches(3.5), fill=LGRAY)
txt(s, "Kontrolleksperiment — XGBoost individuell",
    Inches(6.75), Inches(1.42), Inches(6.1), Inches(0.45),
    size=16, bold=True, color=NAVY)
items_k = [
    {"text": "XGBoost trent separat per SKU", "size": 14, "space": 8},
    {"text": "Samme 36 treningsobservasjoner som HW og ARIMA", "size": 14, "space": 5},
    {"text": "→  ~24 brukbare treningsrader per modell", "size": 13, "italic": True, "space": 3},
    {"text": "", "size": 6},
    {"text": "Resultat:", "bold": True, "size": 15, "color": NAVY, "space": 8},
    {"text": "XGBoost individuell:  MAPE  62,61 %", "size": 15, "bold": True, "color": RED},
    {"text": "Holt-Winters:              MAPE  62,49 %", "size": 15, "bold": True, "color": BLUE, "space": 3},
    {"text": "XGBoost global:          MAPE  52,19 %", "size": 15, "bold": True, "color": GREEN, "space": 3},
]
bullets(s, items_k, Inches(6.75), Inches(1.9), Inches(6.0), Inches(2.8))

rect(s, Inches(0.35), Inches(5.05), Inches(12.6), Inches(1.8), fill=LBLUE)
txt(s, "Konklusjon: Det er global pooling — ikke gradient boosting-arkitekturen — som forklarer XGBoosts fordel.\n"
       "DM-testen gir ikke grunnlag for å fastslå at noen modell er systematisk overlegen.",
    Inches(0.5), Inches(5.12), Inches(12.3), Inches(1.6),
    size=15, bold=True, color=NAVY)

add_notes(s, "Dette er to av de viktigste funnene i studien. DM-testen sier at selv om XGBoost har lavere MAPE enn ARIMA, er ikke forskjellen statistisk signifikant. Og kontrolleksperimentet demonstrerer tydelig at det er pooling (3780 obs) — ikke algoritmen — som gir XGBoosts fordel.")

# ── SLIDE 13: Diskusjon ───────────────────────────────────────────────────────
s = new_slide()
title_bar(s, "Diskusjon — mekanismer bak funnene")
footer(s, 13)

boxes = [
    ("Hvorfor XGBoost svak på sesongpregede SKU-er?",
     "Sesong representeres kun via lagg 12.\n"
     "18/105 SKU-er (17 %) har sterk sesong — en minoritet.\n"
     "Global modell optimerer for de 87 ikke-sesongpregede.\n"
     "Sesongpregede serier «drukner» i støy fra flertallet.\n"
     "HW og ARIMA tilpasser sesong per serie — immune mot dette.",
     LBLUE, NAVY),
    ("Hvorfor XGBoost svak på lavvolum-SKU-er?",
     "Høy CV og nullobservasjoner → laggsignal domineres av støy.\n"
     "Global modell trekker prediksjoner mot porteføljesnittet.\n"
     "Gunstig for høyvolum (ligner snittet),\n"
     "men systematisk skjevhet for lavvolum (avviker fra snittet).",
     LBLUE, NAVY),
    ("Begrensning: asymmetrisk treningsgrunnlag",
     "Global XGBoost: 3 780 obs  vs.  HW/ARIMA: 36 obs/modell.\n"
     "Dette er en bevisst metodisk asymmetri — ikke en feil.\n"
     "Kontrolleksperimentet kvantifiserer og isolerer effekten.\n"
     "Funnet dokumenteres åpent i rapporten.",
     LGRAY, BLACK),
    ("Suksesskriterium MAPE < 30 % — ikke oppnådd",
     "Beste resultat: 52,19 % (XGBoost global).\n"
     "Kriteriet satt før EDA avdekket datakompleksiteten.\n"
     "Burinskiene (2022): 50–70 % er normalt for slike data.\n"
     "Reproduserbarhet og alle fire modeller: fullt oppfylt.",
     LGRAY, BLACK),
]

for i, (head, body, bg, fg) in enumerate(boxes):
    col = i % 2
    row = i // 2
    x = Inches(0.3 + col * 6.5)
    y = Inches(1.35 + row * 2.9)
    rect(s, x, y, Inches(6.3), Inches(2.7), fill=bg)
    txt(s, head, x + Inches(0.15), y + Inches(0.1),
        Inches(6.0), Inches(0.5), size=14, bold=True, color=fg)
    txt(s, body, x + Inches(0.15), y + Inches(0.6),
        Inches(6.0), Inches(2.0), size=12, color=fg if fg != NAVY else BLACK)

add_notes(s, "Forklar mekanismene. To viktige forklaringer: (1) sesong-fortynningseffekten — sesongpregede serier er en minoritet i treningsdata. (2) signal-støy-problemet for lavvolum — den globale modellen passer best for serier som ligner gjennomsnittet.")

# ── SLIDE 14: Konklusjon ──────────────────────────────────────────────────────
s = new_slide()
title_bar(s, "Konklusjon og faglige bidrag")
footer(s, 14)

for i, (num, text, color) in enumerate([
    ("1", "XGBoost global gir observert forbedring over klassiske modeller\n"
          "(MAPE 52,19 % vs. 62,49 % for HW), men DM-testen bekrefter ikke\n"
          "statistisk signifikant overlegenhet over ARIMA (p = 0,436).", BLUE),
    ("2", "Modellprestasjonene varierer systematisk mellom segmenter.\n"
          "XGBoost sterkest for høyvolum og ikke-sesong;\n"
          "ARIMA og Holt-Winters mer robuste for sesongpregede SKU-er.", BLUE),
    ("3", "Ingen av de testede modellparene skiller seg signifikant\n"
          "fra hverandre i DM-test (begge p > 0,05).", BLUE),
    ("4", "Kontrolleksperimentet med XGBoost per SKU (MAPE 62,61 %)\n"
          "isolerer pooling-effekten: global datautnyttelse —\n"
          "ikke modellarkitekturen — forklarer XGBoosts fordel.", BLUE),
]):
    y = Inches(1.4 + i * 1.45)
    rect(s, Inches(0.35), y, Inches(0.7), Inches(1.1), fill=NAVY)
    txt(s, num, Inches(0.35), y + Inches(0.15), Inches(0.7), Inches(0.8),
        size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(1.15), y, Inches(11.8), Inches(1.1), fill=LGRAY)
    txt(s, text, Inches(1.3), y + Inches(0.05),
        Inches(11.5), Inches(1.05), size=14, color=BLACK)

add_notes(s, "Oppsummer de fire faglige bidragene. Understreke at studien er metodisk sterk: felles evalueringsrammeverk, statistisk testing og et kontrolleksperiment som isolerer en spesifikk effekt. Konklusjonen er nyansert — XGBoost er best observert, men ikke statistisk bevisst overlegen ARIMA.")

# ── SLIDE 15: Anbefalinger ────────────────────────────────────────────────────
s = new_slide()
title_bar(s, "Praktiske anbefalinger og videre forskning")
footer(s, 15)

rect(s, Inches(0.35), Inches(1.4), Inches(12.6), Inches(0.4), fill=NAVY)
txt(s, "Praktiske anbefalinger for farmasøytisk distribusjon",
    Inches(0.5), Inches(1.42), Inches(12.3), Inches(0.38),
    size=14, bold=True, color=WHITE)

for icon, head, body, bg, x, w in [
    ("✓", "Med datakompetanse og Python-infrastruktur:",
     "Bruk XGBoost som global modell — gir lavest prognosefeil\nsamlet, særlig for høyvolum-produkter.",
     LGRAY, 0.35, 5.95),
    ("✓", "Uten avansert analysekapasitet:",
     "Holt-Winters er et robust og enkelt alternativ med lav\nimplementeringskostnad og stabil prestasjon.",
     LGRAY, 6.5, 5.95),
]:
    rect(s, Inches(x), Inches(1.9), Inches(w), Inches(1.7), fill=bg)
    txt(s, icon + "  " + head, Inches(x+0.15), Inches(1.97),
        Inches(w-0.3), Inches(0.45), size=13, bold=True, color=NAVY)
    txt(s, body, Inches(x+0.15), Inches(2.45),
        Inches(w-0.3), Inches(1.05), size=13, color=BLACK)

rect(s, Inches(0.35), Inches(3.75), Inches(12.6), Inches(0.4), fill=BLUE)
txt(s, "Videre forskning", Inches(0.5), Inches(3.77),
    Inches(12.3), Inches(0.38), size=14, bold=True, color=WHITE)

videre = [
    {"text": "→  Lengre historikk (72–96 mnd) for å gi XGBoost bedre treningsgrunnlag", "size": 14, "space": 8},
    {"text": "→  Hierarkisk prognosering — aggregere på terapeutisk kategorinivå", "size": 14, "space": 6},
    {"text": "→  Hybrid-strategi: global XGBoost for høyvolum/ikke-sesong + individuelle statistiske modeller for sesong-SKU-er", "size": 14, "space": 6},
    {"text": "→  Inkludere eksogene variabler (epidemiologiske indikatorer, kampanjer)", "size": 14, "space": 6},
    {"text": "→  Globale varianter av HW/ARIMA — sammenligne arkitekturer uten confounding fra ulik datamengde", "size": 14, "space": 6},
]
bullets(s, videre, Inches(0.5), Inches(4.2), Inches(12.5), Inches(2.8))

add_notes(s, "Avslutt med praktiske implikasjoner. For en eksamensituasjon: understreke at anbefalingen avhenger av virksomhetens kapasitet. Nevn at hybrid-strategien er det logiske neste steget basert på segmentanalysen.")

# ── Lagre ─────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Presentasjon lagret: {OUTPUT}")
print(f"Antall slides: {len(prs.slides)}")
